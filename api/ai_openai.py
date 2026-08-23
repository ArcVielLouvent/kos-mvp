# ============================================================
# ai_openai.py -- BACKUP provider OpenAI (bukan dipakai aktif)
# ============================================================
# File ini adalah versi cadangan ai.py yang memakai OpenAI, BUKAN
# Google Gemini. Semua nama fungsi & signature-nya SAMA PERSIS
# dengan ai.py asli, supaya index.py tidak perlu diubah sama sekali
# -- tinggal ganti baris import di index.py dari:
#     from . import ai
# jadi:
#     from . import ai_openai as ai
# dan isi GEMINI_API_KEY di env var diganti/ditambah OPENAI_API_KEY.
#
# CATATAN PENTING sebelum dipakai:
# 1. Belum pernah dites end-to-end (dibuat sebagai jaga-jaga/cadangan,
#    bukan dari hasil debugging nyata seperti ai.py asli).
# 2. Fungsi yang murni olah file lokal (chunk_text, extract_docx_text,
#    extract_pptx_text, extract_xlsx_text, extract_pdf_text_local,
#    extract_pdf_ocr_local, extract_rtf_text, format_dataframe_as_text,
#    is_file_request, is_generate_request, is_analysis_request,
#    infer_doc_type, filter_docs_by_intent, create_docx_bytes,
#    create_pdf_bytes, create_xlsx_bytes, create_docx_from_template)
#    TIDAK provider-specific -- disalin apa adanya dari ai.py, tidak
#    ada logic yang berubah.
# 3. describe_youtube_video() mengambil TRANSKRIP/CAPTION video (lewat
#    youtube-transcript-api) lalu diringkas via GPT -- BUKAN analisis
#    visual seperti Gemini (yang beneran "nonton" videonya). Link video
#    & pemutarannya di UI TIDAK terpengaruh sama sekali oleh ini --
#    fungsi ini cuma menambah teks supaya video lebih mudah DICARI,
#    bukan menyimpan/mengubah link videonya (itu ditangani index.py &
#    komponen SourceLink di frontend, provider-agnostic). Kalau video
#    tidak punya caption sama sekali, tetap return "" dan fallback ke
#    judul+deskripsi manual (perilaku sama seperti sebelumnya).
# 4. extract_multimodal() untuk VIDEO cuma transkrip AUDIO-nya (lewat
#    Whisper) -- TIDAK ada deskripsi visual seperti Gemini. Ini
#    keterbatasan nyata dibanding versi Gemini, bukan bug.
# ============================================================

from concurrent.futures import ThreadPoolExecutor, as_completed
import os
import time
from functools import lru_cache
from openai import OpenAI


def get_client() -> OpenAI:
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise ValueError(
            "Environment variable OPENAI_API_KEY belum di-set!")
    return OpenAI(api_key=api_key)


# ==========================================
# FUNGSI DETEKSI NIAT & UTILITY -- TIDAK provider-specific, disalin
# apa adanya dari ai.py (Gemini), tidak ada perubahan logic.
# ==========================================
def is_file_request(question: str) -> bool:
    keywords = [
        "file asli", "file aslinya", "filenya", "file nya", "download",
        "unduh", "downloadkan", "unduhkan", "kirim file", "kirimkan file",
        "berikan file", "kasih file", "dokumen aslinya", "dokumen asli",
        "dokumennya",
    ]
    q = question.lower()
    return any(kw in q for kw in keywords)


def filter_docs_by_intent(question: str, docs: list) -> list:
    video_keywords = ["video", "youtube", "tonton", "nonton", "putar"]
    q = question.lower()
    if any(kw in q for kw in video_keywords):
        video_docs = [
            d for d in docs if d.get("metadata", {}).get("tipe_file") == "Video YouTube"
        ]
        if video_docs:
            return video_docs
    return docs


def is_generate_request(question: str) -> bool:
    keywords = [
        "buatkan", "buat draf", "buat dokumen", "buatkan dokumen",
        "susun dokumen", "susunkan", "bikinkan", "bikin dokumen",
        "buat form", "buatkan form", "buat formulir", "buatkan formulir",
        "buat sop", "buatkan sop", "buat surat", "buatkan surat",
        "buat kebijakan", "buatkan kebijakan", "generate dokumen",
        "generate draf", "apakah anda bisa membuatkan", "bisa buatkan",
        "bisa dibuatkan",
    ]
    q = question.lower()
    return any(kw in q for kw in keywords)


def infer_doc_type(question: str) -> str:
    q = question.lower()
    if "sop" in q:
        return "SOP"
    if "form" in q or "formulir" in q or "checklist" in q or "check list" in q:
        return "Form/Checklist"
    if "surat" in q:
        return "Surat"
    if "kebijakan" in q:
        return "Kebijakan"
    return "Lainnya"


def is_analysis_request(question: str) -> bool:
    keywords = [
        "analisis", "analisa", "rekomendasi", "rekomendasikan",
        "siapa yang cocok", "cocok jadi", "berdasarkan data",
        "urutkan berdasarkan", "filter data", "dari data yang ada",
    ]
    q = question.lower()
    return any(kw in q for kw in keywords)


def is_compile_request(question: str) -> bool:
    """Deteksi niat: user minta AI MENGGABUNGKAN/MENGEKSTRAK data dari BANYAK
    dokumen berbeda jadi 1 tabel (mis. semua CV di folder, semua laporan
    kinerja) -- BEDA dengan is_analysis_request yang cuma memfilter 1
    dataset XLSX yang sudah terstruktur. Termasuk pertanyaan penilaian/
    rekomendasi lintas-dokumen seperti "siapa yang sebaiknya di-PHK/
    dipromosikan berdasarkan ...".
    """
    keywords = [
        "kompilasi", "kompilasikan", "gabungkan data", "gabungkan semua",
        "rangkum semua", "rekap semua", "rekap dari semua",
        "buatkan daftar", "buatkan tabel", "buat daftar", "buat tabel",
        "bandingkan semua", "dari semua dokumen", "dari semua file",
        "dari semua cv", "dari seluruh dokumen", "ekstrak data dari semua",
        "berdasarkan dokumen-dokumen", "berdasarkan semua dokumen",
        "berdasarkan seluruh laporan", "berdasarkan laporan-laporan",
        "sebaiknya di phk", "sebaiknya dipecat", "sebaiknya diberhentikan",
        "kandidat phk", "layak dipromosikan", "sebaiknya dipromosikan",
        "kandidat promosi", "urutkan karyawan", "ranking karyawan",
        "rangking karyawan",
    ]
    q = question.lower()
    return any(kw in q for kw in keywords)


def synthesize_compiled_answer(question: str, columns: list, rows: list) -> str:
    """Setelah data diekstrak dari banyak dokumen jadi tabel (rows), AI
    menjawab pertanyaan ASLI user berdasarkan tabel itu SAJA -- termasuk
    pertanyaan penilaian/rekomendasi (PHK, promosi, dst). Jawaban WAJIB
    berpijak ke baris data yang ada (sebut nama/baris mana yang mendukung
    kesimpulan), bukan menebak di luar data.

    Untuk pertanyaan yang menyangkut keputusan personel (PHK/promosi/dst),
    jawaban tetap diminta menyertakan alasan dari data, tapi ditutup
    dengan pengingat singkat bahwa ini bantuan awal berbasis dokumen yang
    tersedia -- keputusan final tetap perlu pertimbangan HR/atasan
    (bukan AI) dan konteks yang mungkin tidak tercatat di dokumen.
    """
    table_text = " | ".join(columns) + "\n"
    for r in rows:
        table_text += " | ".join(str(r.get(c, "-")) for c in columns) + "\n"

    prompt = f"""Berikut tabel hasil ekstraksi data dari {len(rows)} dokumen:

{table_text}

Pertanyaan user: "{question}"

Jawab pertanyaan itu HANYA berdasarkan tabel di atas. Sebutkan nama/baris spesifik yang
mendukung jawabanmu. Kalau data di tabel tidak cukup untuk menjawab dengan yakin, katakan
itu terus terang -- JANGAN menebak atau menambah informasi yang tidak ada di tabel.

Kalau pertanyaan menyangkut keputusan personel (PHK, promosi, mutasi, dst), tutup jawabanmu
dengan 1 kalimat singkat bahwa ini bantuan awal berbasis dokumen yang tersedia, dan keputusan
akhir tetap perlu pertimbangan manusia (HR/atasan) karena mungkin ada konteks yang tidak
tercatat di dokumen.
"""
    return (_generate_with_fallback([{"role": "user", "content": prompt}]) or "").strip()


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list:
    if not text:
        return []
    chunks = []
    i = 0
    while i < len(text):
        end = i + chunk_size
        if end < len(text):
            space_index = text.rfind(" ", i, end)
            if space_index > i:
                end = space_index
        chunks.append(text[i:end].strip())
        i = end - overlap
    return [c for c in chunks if len(c) > 10]


# ==========================================
# RETRY HELPER (setara _call_with_retry di ai.py)
# ==========================================
def _call_with_retry(func, *args, max_retries: int = 4, base_delay: int = 3, **kwargs):
    last_error = None
    for attempt in range(max_retries):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            last_error = e
            error_str = str(e)
            is_transient = any(
                code in error_str
                for code in ["503", "500", "429", "rate_limit", "RESOURCE_EXHAUSTED", "server_error"]
            )
            if is_transient and attempt < max_retries - 1:
                wait = base_delay * (2 ** attempt)
                time.sleep(wait)
                continue
            raise
    raise last_error


# Urutan: model murah/cepat dulu, fallback ke yang lebih lawas/stabil
# kalau ada masalah kuota/akses. SESUAIKAN kalau OpenAI merilis model
# baru -- ini bukan daftar permanen, cek platform.openai.com/docs/models
CHAT_FALLBACK_MODELS = ["gpt-5-mini", "gpt-4o-mini"]
EMBEDDING_FALLBACK_MODELS = ["text-embedding-3-small", "text-embedding-3-large"]
EMBEDDING_DIMENSIONS = 768  # samain dengan kolom vector(768) di schema


def _generate_with_fallback(messages: list) -> str:
    """Setara _generate_with_fallback di ai.py, tapi pakai Chat Completions API."""
    client = get_client()
    last_error = None

    for model_name in CHAT_FALLBACK_MODELS:
        try:
            response = _call_with_retry(
                client.chat.completions.create,
                model=model_name,
                messages=messages,
            )
            return response.choices[0].message.content or ""
        except Exception as e:
            last_error = e
            error_str = str(e)
            if any(code in error_str for code in ["429", "rate_limit", "insufficient_quota", "403", "404", "model_not_found"]):
                continue
            raise

    raise last_error


def embed_text(text: str) -> list:
    """Setara embed_text di ai.py -- dimensions dikunci 768 biar cocok
    dengan schema database yang sudah ada, gak perlu migrasi kolom."""
    client = get_client()
    last_error = None

    for model_name in EMBEDDING_FALLBACK_MODELS:
        try:
            result = _call_with_retry(
                client.embeddings.create,
                model=model_name,
                input=text,
                dimensions=EMBEDDING_DIMENSIONS,
            )
            print(f"[DEBUG-EMBED-OPENAI] Sukses pakai model '{model_name}'")
            return result.data[0].embedding
        except Exception as e:
            last_error = e
            error_str = str(e)
            print(f"[DEBUG-EMBED-OPENAI] Gagal model '{model_name}': {error_str}")
            if any(code in error_str for code in ["429", "rate_limit", "insufficient_quota", "403", "404", "model_not_found"]):
                continue
            raise

    raise last_error


def embed_chunks_parallel(chunks: list, max_workers: int = 4) -> list:
    results = [None] * len(chunks)
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_index = {
            executor.submit(embed_text, chunk): i for i, chunk in enumerate(chunks)
        }
        for future in as_completed(future_to_index):
            i = future_to_index[future]
            results[i] = future.result()
    return results


def generate_answer(question: str, context_documents: list) -> str:
    context_parts = []
    for doc in context_documents:
        sumber = doc.get("metadata", {}).get("tipe_file", "Dokumen KOS")
        folder = doc.get("folder_path", "/")
        context_parts.append(
            f"Judul: {doc['title']}\nLokasi: {folder}\nTipe: {sumber}\nIsi:\n{doc['content'][:3000]}"
        )
    context = "\n\n====================\n\n".join(context_parts)

    system_prompt = """Kamu adalah Knowledge Operating System (KOS), asisten cerdas internal perusahaan.
Jawab pertanyaan pengguna HANYA berdasarkan Dokumen Referensi yang diberikan.
Jika jawaban tidak ada di dalam dokumen referensi, katakan jujur bahwa informasi tersebut belum tersedia di database. Jangan pernah mengarang.

ATURAN KHUSUS UNTUK DATA TABEL (baris berformat "kolom1 | kolom2 | ..."):
- Jika sumber data berupa tabel/sheet, tampilkan sebagai tabel Markdown (bisa di-copy),
  BUKAN diringkas atau diparafrasekan dengan kalimat bebas.
- Salin nilai apa adanya persis seperti di sumber -- jangan mengubah, membulatkan, atau menerka angka."""

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"=== DOKUMEN REFERENSI ===\n{context}\n=== AKHIR DOKUMEN ===\n\nPertanyaan: {question}"},
    ]
    return _generate_with_fallback(messages)


def extract_multimodal(file_path: str, mime_type: str, display_name: str) -> str:
    """
    KETERBATASAN dibanding versi Gemini:
    - PDF & gambar: pakai Vision (gpt model bisa baca gambar) -- OK, setara.
    - Audio: pakai Whisper -- OK, setara.
    - VIDEO: OpenAI TIDAK bisa "menonton" video. Fungsi ini cuma transkrip
      AUDIO dari video-nya (bukan deskripsi visual). Ini keterbatasan nyata,
      bukan bug -- kalau butuh deskripsi visual video, provider Gemini
      masih lebih unggul untuk kasus ini.
    """
    client = get_client()

    try:
        if "pdf" in mime_type or "image" in mime_type:
            import base64
            with open(file_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("utf-8")

            if "pdf" in mime_type:
                prompt_text = (
                    "Baca seluruh dokumen PDF ini dengan saksama dan ekstrak seluruh "
                    "teks serta tabel menjadi teks terstruktur murni."
                )
            else:
                prompt_text = (
                    "Amati gambar ini dengan saksama. Jika ada teks di dalamnya, "
                    "transkrip teksnya secara lengkap dan akurat. Jika ini foto biasa "
                    "tanpa teks, deskripsikan isinya secara detail dan faktual."
                )

            messages = [{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt_text},
                    {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{b64}"}},
                ],
            }]
            hasil_teks = _generate_with_fallback(messages)

        elif "audio" in mime_type or "video" in mime_type:
            # Video: cuma audio-nya yang ditranskrip (lihat catatan di docstring)
            with open(file_path, "rb") as f:
                transcript = _call_with_retry(
                    client.audio.transcriptions.create,
                    model="whisper-1",
                    file=f,
                )
            hasil_teks = transcript.text
        else:
            raise RuntimeError(f"Tipe file '{mime_type}' belum didukung oleh backend OpenAI.")

        if not hasil_teks or not hasil_teks.strip():
            raise ValueError(f"Tidak ada teks yang bisa diekstrak dari '{display_name}'.")

        return hasil_teks

    except Exception as e:
        raise RuntimeError(
            f"Gagal mengekstrak '{display_name}' via OpenAI: {str(e)}"
        )


# ==========================================
# EKSTRAKSI LOKAL (DOCX / PPTX / XLSX / PDF / RTF)
# Tidak provider-specific -- disalin apa adanya dari ai.py.
# ==========================================
def extract_docx_text(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def extract_pptx_text(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        texts.append(text)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_xlsx_text(file_path: str) -> list:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join(str(c) if c is not None else "" for c in row))
        if rows:
            sheets.append((sheet.title, "\n".join(rows)))
    return sheets


def extract_xlsx_structured(file_path: str) -> list:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    result = []
    for sheet in wb.worksheets:
        all_rows = list(sheet.iter_rows(values_only=True))
        if len(all_rows) < 2:
            continue
        header = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(all_rows[0])]
        rows = []
        for row in all_rows[1:]:
            if any(cell is not None for cell in row):
                rows.append({header[i]: row[i] for i in range(min(len(header), len(row)))})
        if rows:
            result.append({"sheet": sheet.title, "rows": rows})
    return result


def format_dataframe_as_text(df, sheet_name: str = "Data") -> str:
    header = " | ".join(str(c) for c in df.columns)
    rows = [" | ".join(str(v) for v in row) for row in df.itertuples(index=False)]
    return f"Sheet: {sheet_name}\n" + header + "\n" + "\n".join(rows)


def describe_youtube_video(youtube_url: str) -> str:
    """
    OpenAI tidak bisa 'menonton' video langsung seperti Gemini File API,
    tapi bisa ambil TRANSKRIP/CAPTION video (kalau tersedia) lalu
    diringkas pakai GPT. Ini analisis berbasis TEKS (caption), bukan
    visual -- kalau video minim narasi, hasilnya kurang lengkap
    dibanding versi Gemini. Return "" kalau caption tidak tersedia
    sama sekali (caller sudah fallback ke judul+deskripsi manual).
    """
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        import re

        match = re.search(r"(?:v=|youtu\.be/|embed/|shorts/)([a-zA-Z0-9_-]{11})", youtube_url)
        if not match:
            return ""
        video_id = match.group(1)

        transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=["id", "en"])
        full_text = " ".join(seg["text"] for seg in transcript_list)

        if not full_text.strip():
            return ""

        prompt = (
            "Berikut adalah transkrip sebuah video. Buatkan ringkasan konten yang detail: "
            "topik utama, langkah-langkah atau poin penting yang dibahas, dan konteks lain "
            "yang relevan untuk pencarian internal perusahaan.\n\n"
            f"=== TRANSKRIP ===\n{full_text[:8000]}\n=== AKHIR TRANSKRIP ==="
        )
        return _generate_with_fallback([{"role": "user", "content": prompt}]) or ""
    except Exception:
        return ""


def extract_pdf_text_local(file_path: str) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            if text.strip():
                parts.append(f"Halaman {i + 1}:\n{text}")
    return "\n\n".join(parts)


def extract_pdf_ocr_local(file_path: str) -> str:
    import pytesseract
    from pdf2image import convert_from_path
    images = convert_from_path(file_path)
    parts = []
    for i, image in enumerate(images):
        text = pytesseract.image_to_string(image, lang="ind+eng")
        if text.strip():
            parts.append(f"Halaman {i + 1}:\n{text}")
    return "\n\n".join(parts)


def extract_rtf_text(file_path: str) -> str:
    from striprtf.striprtf import rtf_to_text
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    return rtf_to_text(raw)


# ==========================================
# GENERATE KUIS TRAINING
# ==========================================
def generate_quiz_questions(content: str, num_questions: int = 5) -> list:
    prompt = f"""Kamu adalah pembuat soal ujian untuk training karyawan baru.
Berdasarkan dokumen di bawah, buatkan {num_questions} soal pilihan ganda (4 opsi jawaban)
yang menguji pemahaman terhadap SOP/prosedur di dokumen ini.

ATURAN KETAT:
- Jawab HANYA dengan JSON murni, tanpa teks penjelasan apa pun di luar JSON.
- Jangan pakai markdown code block (tanpa ```json).
- Format persis seperti ini:
[
  {{"question": "...", "options": ["A", "B", "C", "D"], "correct_index": 0}},
  ...
]
- correct_index adalah index (0-3) dari jawaban yang benar di array "options".
- Soal harus berdasarkan isi dokumen, bukan pengetahuan umum di luar dokumen.

=== ISI DOKUMEN ===
{content[:6000]}
=== AKHIR DOKUMEN ===
"""
    raw = (_generate_with_fallback([{"role": "user", "content": prompt}]) or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()

    import json
    questions = json.loads(raw)

    valid = []
    for q in questions:
        if (
            isinstance(q, dict) and "question" in q and "options" in q
            and len(q["options"]) == 4 and "correct_index" in q
            and 0 <= q["correct_index"] <= 3
        ):
            valid.append(q)
    return valid


# ==========================================
# GENERATE DOKUMEN AI
# ==========================================
DOC_TYPE_INSTRUCTIONS = {
    "SOP": "Format sebagai SOP: tujuan, ruang lingkup, lalu langkah-langkah bernomor yang jelas dan actionable.",
    "Form/Checklist": "Format sebagai daftar item/kolom isian singkat, cocok dijadikan tabel formulir.",
    "Surat": "Format sebagai surat resmi: kop (placeholder), tanggal, salam pembuka, isi, salam penutup, tempat tanda tangan.",
    "Kebijakan": "Format sebagai pernyataan kebijakan: latar belakang, ketentuan, sanksi/konsekuensi kalau relevan.",
    "Lainnya": "Format bebas namun tetap terstruktur dengan heading yang jelas.",
}


def generate_draft_document(topic: str, doc_type: str = "Lainnya", company_context: str = "") -> str:
    instruction = DOC_TYPE_INSTRUCTIONS.get(doc_type, DOC_TYPE_INSTRUCTIONS["Lainnya"])
    prompt = f"""Kamu diminta membuat DRAF dokumen kerja untuk sebuah perusahaan.

Jenis dokumen: {doc_type}
Topik: {topic}
Konteks perusahaan (kalau ada): {company_context or "Tidak ada konteks tambahan."}

Instruksi format: {instruction}

ATURAN WAJIB:
- Di baris PALING ATAS, tulis persis: "[DRAF AI -- PERLU DIREVIEW SEBELUM DIGUNAKAN RESMI]"
- JANGAN mengarang detail teknis spesifik yang tidak bisa dipastikan kebenarannya --
  gunakan placeholder seperti "[isi sesuai SOP internal]".
"""
    return _generate_with_fallback([{"role": "user", "content": prompt}]) or ""


def create_docx_bytes(title: str, content: str, logo_bytes: bytes = None) -> bytes:
    from docx import Document
    from docx.shared import Inches
    import io

    doc = Document()
    if logo_bytes:
        try:
            doc.add_picture(io.BytesIO(logo_bytes), width=Inches(1.2))
        except Exception:
            pass
    doc.add_heading(title, level=1)
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("[DRAF AI"):
            p = doc.add_paragraph()
            run = p.add_run(line)
            run.bold = True
        elif line.startswith("#"):
            doc.add_heading(line.lstrip("#").strip(), level=2)
        else:
            doc.add_paragraph(line)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def create_docx_from_template(template_bytes: bytes, title: str, content: str) -> bytes:
    from docx import Document
    import io

    doc = Document(io.BytesIO(template_bytes))
    doc.add_heading(title, level=1)
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("[DRAF AI"):
            p = doc.add_paragraph()
            run = p.add_run(line)
            run.bold = True
        elif line.startswith("#"):
            doc.add_heading(line.lstrip("#").strip(), level=2)
        else:
            doc.add_paragraph(line)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def create_pdf_bytes(title: str, content: str, logo_bytes: bytes = None) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    import io

    styles = getSampleStyleSheet()
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm,
        leftMargin=2 * cm, rightMargin=2 * cm,
    )
    story = []
    if logo_bytes:
        try:
            story.append(Image(io.BytesIO(logo_bytes), width=3 * cm, height=3 * cm))
            story.append(Spacer(1, 12))
        except Exception:
            pass
    story.append(Paragraph(title, styles["Title"]))
    story.append(Spacer(1, 16))
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 8))
            continue
        style = (
            styles["Heading2"] if line.startswith("[DRAF AI") or line.startswith("#") else styles["Normal"]
        )
        story.append(Paragraph(line.lstrip("#").strip(), style))
    doc.build(story)
    return buffer.getvalue()


def create_xlsx_bytes(title: str, rows: list) -> bytes:
    import openpyxl
    import io

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31] if title else "Data"
    if rows:
        headers = list(rows[0].keys())
        ws.append(headers)
        for row in rows:
            ws.append([row.get(h, "") for h in headers])

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


# ==========================================
# ANALISIS DATA
# ==========================================
def determine_extraction_columns(user_request: str) -> list:
    prompt = f"""Permintaan pengguna: "{user_request}"

Tentukan kolom-kolom apa saja yang seharusnya ada di tabel hasil kompilasi data.
Jawab HANYA dengan JSON array of string, tanpa markdown code block, contoh:
["Nama", "Posisi Dilamar", "Pengalaman (tahun)"]
"""
    raw = (_generate_with_fallback([{"role": "user", "content": prompt}]) or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    import json
    return json.loads(raw)


def extract_fields_from_document(content: str, columns: list, source_title: str) -> dict:
    prompt = f"""Dokumen sumber ("{source_title}"):
{content[:4000]}

Ekstrak nilai untuk kolom berikut PERSIS dari isi dokumen di atas: {columns}
Kalau suatu info tidak disebutkan di dokumen, isi nilainya dengan "-" -- JANGAN menebak/mengarang.

Jawab HANYA dengan JSON object, tanpa markdown code block, contoh:
{{"Nama": "...", "Posisi Dilamar": "..."}}
"""
    raw = (_generate_with_fallback([{"role": "user", "content": prompt}]) or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    import json
    return json.loads(raw)


def extract_fields_from_documents_parallel(
    documents: list, columns: list, max_workers: int = 4
) -> tuple:
    """Jalankan extract_fields_from_document ke BANYAK dokumen sekaligus
    secara paralel -- dipakai fitur kompilasi data lintas-dokumen di Chat
    KOS supaya tidak lambat kalau folder-nya berisi belasan/puluhan
    dokumen. `documents` = list of {"id", "title", "content"}. Dokumen
    yang gagal diekstrak DILEWATI, bukan menggagalkan semuanya -- error
    dikumpulkan terpisah. Return (rows, errors)."""
    results = [None] * len(documents)
    errors = []

    def _extract_one(doc):
        row = extract_fields_from_document(doc["content"], columns, doc["title"])
        row["_source_title"] = doc["title"]
        row["_source_id"] = doc["id"]
        return row

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_index = {
            executor.submit(_extract_one, doc): i for i, doc in enumerate(documents)
        }
        for future in as_completed(future_to_index):
            i = future_to_index[future]
            try:
                results[i] = future.result()
            except Exception as e:
                errors.append(f"{documents[i]['title']}: {str(e)}")

    rows = [r for r in results if r is not None]
    return rows, errors


def extract_analysis_criteria(question: str, columns: list) -> dict:
    prompt = f"""Kolom yang tersedia di data: {columns}

Pertanyaan pengguna: "{question}"

Terjemahkan pertanyaan ini jadi kriteria filter terstruktur. Jawab HANYA dengan JSON murni,
tanpa markdown code block, format persis:
{{
  "missing_info": null,
  "filters": [{{"column": "nama_kolom", "operator": ">=", "value": "1"}}],
  "sort_by": null,
  "sort_desc": false
}}

Operator yang valid HANYA: ">=", "<=", "==", "contains".
Kalau pertanyaan terlalu ambigu / kriteria penting belum disebutkan, isi "missing_info" dengan
pertanyaan klarifikasi singkat dalam Bahasa Indonesia, dan biarkan "filters" jadi array kosong.
"""
    raw = (_generate_with_fallback([{"role": "user", "content": prompt}]) or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    import json
    return json.loads(raw)
