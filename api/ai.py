from concurrent.futures import ThreadPoolExecutor, as_completed
import os
import time
from functools import lru_cache
from google import genai
from google.genai import types


def is_file_request(question: str) -> bool:
    """
    Deteksi niat sederhana: apakah user minta FILE ASLI (bukan jawaban teks)?
    Keyword matching -- gratis & instan, cukup untuk kasus umum. Kalau nanti
    banyak frasa lolos tak terdeteksi, baru pertimbangkan upgrade ke klasifikasi
    berbasis AI (butuh 1 API call tambahan, ada biaya & latency).
    """
    keywords = [
        "file asli",
        "file aslinya",
        "filenya",
        "file nya",
        "download",
        "unduh",
        "downloadkan",
        "unduhkan",
        "kirim file",
        "kirimkan file",
        "berikan file",
        "kasih file",
        "dokumen aslinya",
        "dokumen asli",
        "dokumennya",
    ]
    q = question.lower()
    return any(kw in q for kw in keywords)


def filter_docs_by_intent(question: str, docs: list) -> list:
    """
    Kalau user secara eksplisit minta VIDEO, saring hasil pencarian supaya
    cuma dokumen bertipe Video YouTube yang ditampilkan -- mencegah dokumen
    lain (xlsx/pdf) yang kebetulan mirip kata kunci ikut nongol sebagai
    jawaban/tombol yang tidak relevan.
    """
    video_keywords = ["video", "youtube", "tonton", "nonton", "putar"]
    q = question.lower()
    if any(kw in q for kw in video_keywords):
        video_docs = [
            d for d in docs if d.get("metadata", {}).get("tipe_file") == "Video YouTube"
        ]
        if video_docs:
            return video_docs
    return docs


def is_generate_request(question: str) -> bool:
    """Deteksi niat: user minta dokumen/form/surat/SOP DIBUATKAN, bukan dicari."""
    keywords = [
        "buatkan",
        "buat draf",
        "buat dokumen",
        "buatkan dokumen",
        "susun dokumen",
        "susunkan",
        "bikinkan",
        "bikin dokumen",
        "buat form",
        "buatkan form",
        "buat formulir",
        "buatkan formulir",
        "buat sop",
        "buatkan sop",
        "buat surat",
        "buatkan surat",
        "buat kebijakan",
        "buatkan kebijakan",
        "generate dokumen",
        "generate draf",
        "apakah anda bisa membuatkan",
        "bisa buatkan",
        "bisa dibuatkan",
    ]
    q = question.lower()
    return any(kw in q for kw in keywords)


def infer_doc_type(question: str) -> str:
    """Tebak jenis dokumen dari kata kunci di pertanyaan -- default 'Lainnya'."""
    q = question.lower()
    if "sop" in q:
        return "SOP"
    if "form" in q or "formulir" in q or "checklist" in q or "check list" in q:
        return "Form/Checklist"
    if "surat" in q:
        return "Surat"
    if "kebijakan" in q:
        return "Kebijakan"
    return "Lainnya"


def is_analysis_request(question: str) -> bool:
    """Deteksi niat: user minta analisis/rekomendasi dari data terstruktur (xlsx)."""
    keywords = [
        "analisis",
        "analisa",
        "rekomendasi",
        "rekomendasikan",
        "siapa yang cocok",
        "cocok jadi",
        "berdasarkan data",
        "urutkan berdasarkan",
        "filter data",
        "dari data yang ada",
    ]
    q = question.lower()
    return any(kw in q for kw in keywords)


def embed_chunks_parallel(chunks: list, max_workers: int = 4) -> list:
    """
    Embed banyak chunk SEKALIGUS secara paralel (bukan satu-satu berurutan).
    Mempercepat upload file besar/banyak chunk secara signifikan.
    max_workers dibatasi (bukan tanpa batas) supaya tidak memicu rate limit --
    kalau tetap kena, _call_with_retry di embed_text yang menangani otomatis.
    """
    results = [None] * len(chunks)
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_index = {
            executor.submit(embed_text, chunk): i for i, chunk in enumerate(chunks)
        }
        for future in as_completed(future_to_index):
            i = future_to_index[future]
            results[i] = future.result()
    return results


def get_client() -> genai.Client:
    """Inisialisasi Client menggunakan SDK google-genai yang mutakhir"""
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise ValueError(
            "Environment variable GEMINI_API_KEY belum di-set di Vercel!")
    return genai.Client(api_key=api_key)


def _call_with_retry(func, *args, max_retries: int = 4, base_delay: int = 3, **kwargs):
    """
    Coba ulang otomatis kalau kena error transient dari server Google
    (503 UNAVAILABLE / 429 rate limit) -- bukan bug di kode, murni server
    Google sedang sibuk. Backoff: 3s, 6s, 12s, 24s.
    """
    last_error = None
    for attempt in range(max_retries):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            last_error = e
            error_str = str(e)
            is_transient = any(
                code in error_str
                for code in ["503", "UNAVAILABLE", "429", "RESOURCE_EXHAUSTED"]
            )
            if is_transient and attempt < max_retries - 1:
                wait = base_delay * (2**attempt)
                time.sleep(wait)
                continue
            raise
    raise last_error


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list:
    """Sistem pemotong teks (Chunking) untuk mencegah Limit Token"""
    if not text:
        return []
    chunks = []
    i = 0
    while i < len(text):
        end = i + chunk_size
        if end < len(text):
            space_index = text.rfind(" ", i, end)
            if space_index > i:
                end = space_index
        chunks.append(text[i:end].strip())
        i = end - overlap
    return [c for c in chunks if len(c) > 10]


FALLBACK_MODELS = ["gemini-3.5-flash", "gemini-2.5-flash", "gemini-2.0-flash"]


def _generate_with_fallback(contents):
    """
    Coba generate_content dengan beberapa model Gemini berurutan.
    - Kalau kena KUOTA HABIS (429 RESOURCE_EXHAUSTED) -> langsung lompat ke
      model berikutnya di FALLBACK_MODELS (percuma ditunggu, kuota per-model beda).
    - Kalau kena error transient (503/dsb) -> tetap coba ulang di model yang sama
      dulu lewat _call_with_retry, baru lompat model kalau itu juga gagal terus.
    """
    client = get_client()
    last_error = None

    for model_name in FALLBACK_MODELS:
        try:
            return _call_with_retry(
                client.models.generate_content,
                model=model_name,
                contents=contents,
            )
        except Exception as e:
            last_error = e
            error_str = str(e)
            if "RESOURCE_EXHAUSTED" in error_str or "429" in error_str:
                continue  # kuota model ini habis -> lanjut ke model berikutnya
            raise  # error lain (bukan soal kuota), jangan asal ganti model

    raise last_error


@lru_cache(maxsize=1)
def get_embedding_model() -> str:
    """Otomatis beralih ke lini Gemini Embedding terbaru"""
    return "gemini-embedding-2"


@lru_cache(maxsize=1)
def get_generation_model() -> str:
    """Menggunakan lini Gemini Flash paling mutakhir dan stabil"""
    return "gemini-3.5-flash"


EMBEDDING_FALLBACK_MODELS = ["gemini-embedding-001", "gemini-embedding-2"]


def embed_text(text: str) -> list:
    """
    Sintaks ekstraksi array embedding dengan pembatasan dimensi ke 768.
    Fallback ke model embedding berikutnya kalau model saat ini kena kuota
    habis (429) ATAU tidak bisa diakses (403/404) -- bukan cuma kuota.
    """
    client = get_client()
    last_error = None

    for model_name in EMBEDDING_FALLBACK_MODELS:
        try:
            result = _call_with_retry(
                client.models.embed_content,
                model=model_name,
                contents=text,
                config=types.EmbedContentConfig(output_dimensionality=768),
            )
            [embedding_obj] = result.embeddings
            print(f"[DEBUG-EMBED] Sukses pakai model '{model_name}'")
            return embedding_obj.values
        except Exception as e:
            last_error = e
            error_str = str(e)
            print(f"[DEBUG-EMBED] Gagal model '{model_name}': {error_str}")
            is_model_unavailable = any(
                code in error_str
                for code in ["RESOURCE_EXHAUSTED", "429", "PERMISSION_DENIED", "403", "NOT_FOUND", "404"]
            )
            if is_model_unavailable:
                continue
            raise

    raise last_error


def generate_answer(question: str, context_documents: list) -> str:
    client = get_client()
    context_parts = []
    for doc in context_documents:
        sumber = doc.get("metadata", {}).get("tipe_file", "Dokumen KOS")
        folder = doc.get("folder_path", "/")
        context_parts.append(
            f"Judul: {doc['title']}\nLokasi: {folder}\nTipe: {sumber}\nIsi:\n{doc['content']}"
        )

    context = "\n\n====================\n\n".join(context_parts)

    prompt = f"""Kamu adalah Knowledge Operating System (KOS), asisten cerdas internal perusahaan.
Jawab pertanyaan pengguna HANYA berdasarkan Dokumen Referensi di bawah ini.
Jika jawaban tidak ada di dalam dokumen referensi, katakan jujur bahwa informasi tersebut belum tersedia di database. Jangan pernah mengarang.

ATURAN KHUSUS UNTUK DATA TABEL (baris berformat "kolom1 | kolom2 | ..."):
- Jika sumber data berupa tabel/sheet, tampilkan sebagai tabel Markdown (bisa di-copy),
  BUKAN diringkas atau diparafrasekan dengan kalimat bebas.
- Salin nilai apa adanya persis seperti di sumber -- jangan mengubah, membulatkan, atau menerka angka.
- Jika pengguna hanya bertanya sebagian (misal satu baris/kategori tertentu), tampilkan baris relevan
  saja, tapi tetap dalam format tabel.
- Jika pengguna eksplisit minta "seluruh data" / "semua isi sheet ini", tampilkan seluruh baris
  dari sheet yang relevan itu secara lengkap.

=== DOKUMEN REFERENSI ===
{context}
=== AKHIR DOKUMEN ===

Pertanyaan: {question}
"""
    response = _generate_with_fallback(prompt)
    return response.text


def extract_multimodal(file_path: str, mime_type: str, display_name: str) -> str:
    """
    Ekstraksi PDF, audio, dan video via Gemini File API (SDK google-genai).
    File diupload ke Google Cloud, dipoll sampai status ACTIVE (anti-halusinasi),
    lalu dihapus dari server setelah selesai.
    """
    client = get_client()

    try:
        uploaded_file = _call_with_retry(client.files.upload, file=file_path)

        # Polling anti-halusinasi: tunggu sampai file benar-benar siap diproses
        while uploaded_file.state.name == "PROCESSING":
            time.sleep(3)
            uploaded_file = client.files.get(name=uploaded_file.name)

        if uploaded_file.state.name == "FAILED":
            try:
                client.files.delete(name=uploaded_file.name)
            except Exception:
                pass
            raise ValueError(
                f"Google Gemini gagal memproses file '{display_name}'.")

        if "pdf" in mime_type:
            prompt = (
                "Baca seluruh dokumen PDF ini dengan saksama dan ekstrak seluruh "
                "teks serta tabel menjadi teks terstruktur murni."
            )
        elif "video" in mime_type:
            prompt = (
                "Tonton video ini dari awal hingga akhir. Buatkan transkrip sangat detail, "
                "termasuk kejadian visual, langkah-langkah teknis yang ditunjukkan, "
                "dan teks apa pun yang muncul di layar."
            )
        elif "audio" in mime_type:
            prompt = "Dengarkan audio ini dengan saksama dan buatkan transkrip teks yang lengkap dan akurat."
        elif "image" in mime_type:
            prompt = (
                "Amati gambar ini dengan saksama. Jika ada teks di dalamnya (dokumen yang "
                "difoto, poster, papan tulis, tangkapan layar), transkrip teksnya secara "
                "lengkap dan akurat. Jika ini foto biasa tanpa teks, deskripsikan isinya "
                "secara detail dan faktual."
            )
        else:
            prompt = "Ekstrak seluruh informasi dari file ini menjadi teks terstruktur murni."

        response = _generate_with_fallback([uploaded_file, prompt])

        try:
            client.files.delete(name=uploaded_file.name)
        except Exception:
            pass

        hasil_teks = response.text
        if not hasil_teks or not hasil_teks.strip():
            raise ValueError(
                f"Tidak ada teks yang bisa diekstrak dari '{display_name}'."
            )

        return hasil_teks

    except Exception as e:
        raise RuntimeError(
            f"Gagal mengekstrak '{display_name}' via Google File API: {str(e)}"
        )


# ==========================================
# EKSTRAKSI LOKAL (DOCX / PPTX / XLSX)
# Tidak lewat Gemini -- dokumen office tidak dipahami baik oleh document vision
# ==========================================
def extract_docx_text(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    return "\n".join(parts)


def extract_pptx_text(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        texts.append(text)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))

    return "\n\n".join(slides)


def extract_xlsx_text(file_path: str) -> list:
    """
    Kembalikan list per-sheet: [(nama_sheet, isi_lengkap_sheet), ...]
    Sengaja TIDAK digabung jadi satu teks panjang -- supaya tiap sheet
    tetap utuh sebagai satu chunk, tidak terpotong sembarang karakter.
    """
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    sheets = []

    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join(
                    str(c) if c is not None else "" for c in row))
        if rows:
            sheets.append((sheet.title, "\n".join(rows)))

    return sheets


def extract_xlsx_structured(file_path: str) -> list:
    """
    Kembalikan data XLSX dalam bentuk terstruktur untuk kebutuhan analisis
    (bukan cuma teks). Baris pertama tiap sheet dianggap header/nama kolom.
    Format: [{"sheet": "nama", "rows": [{"kolom1": val, ...}, ...]}, ...]
    """
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    result = []

    for sheet in wb.worksheets:
        all_rows = list(sheet.iter_rows(values_only=True))
        if len(all_rows) < 2:
            continue
        header = [
            str(h) if h is not None else f"col_{i}" for i, h in enumerate(all_rows[0])
        ]
        rows = []
        for row in all_rows[1:]:
            if any(cell is not None for cell in row):
                rows.append(
                    {header[i]: row[i]
                        for i in range(min(len(header), len(row)))}
                )
        if rows:
            result.append({"sheet": sheet.title, "rows": rows})

    return result


def format_dataframe_as_text(df, sheet_name: str = "Data") -> str:
    """Ubah pandas DataFrame (mis. dari CSV) jadi teks tabel utuh, tanpa dipotong per baris."""
    header = " | ".join(str(c) for c in df.columns)
    rows = [" | ".join(str(v) for v in row)
            for row in df.itertuples(index=False)]
    return f"Sheet: {sheet_name}\n" + header + "\n" + "\n".join(rows)


def describe_youtube_video(youtube_url: str) -> str:
    """
    Minta Gemini menonton video YouTube LANGSUNG dari URL-nya (tanpa perlu
    download/upload file). Best-effort: kalau gagal (fitur berubah, video
    private, dsb), kembalikan string kosong -- pemanggil tinggal fallback
    ke judul+deskripsi manual saja.
    """
    try:
        prompt = (
            "Tonton video ini dan buatkan ringkasan konten yang detail: topik utama, "
            "langkah-langkah atau poin penting yang dibahas, dan konteks lain yang "
            "relevan untuk pencarian internal perusahaan."
        )
        contents = types.Content(
            parts=[
                types.Part(file_data=types.FileData(file_uri=youtube_url)),
                types.Part(text=prompt),
            ]
        )
        response = _generate_with_fallback(contents)
        return response.text or ""
    except Exception:
        return ""


def extract_pdf_text_local(file_path: str) -> str:
    """
    Ekstraksi teks PDF LOKAL (pdfplumber) -- gratis, instan, TIDAK memakai kuota
    Gemini sama sekali. Cocok untuk PDF teks digital biasa (bukan hasil scan).
    Return string kosong kalau PDF-nya scan/gambar (tidak ada teks terbaca).
    """
    import pdfplumber

    parts = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            if text.strip():
                parts.append(f"Halaman {i + 1}:\n{text}")
    return "\n\n".join(parts)


def extract_pdf_ocr_local(file_path: str) -> str:
    """
    OCR lokal pakai Tesseract -- gratis, tanpa API, tanpa kuota Gemini sama sekali.
    Dipakai khusus untuk PDF hasil scan/gambar yang tidak punya teks digital
    (extract_pdf_text_local akan kosong untuk PDF jenis ini).
    """
    import pytesseract
    from pdf2image import convert_from_path

    images = convert_from_path(file_path)
    parts = []
    for i, image in enumerate(images):
        text = pytesseract.image_to_string(image, lang="ind+eng")
        if text.strip():
            parts.append(f"Halaman {i + 1}:\n{text}")
    return "\n\n".join(parts)


def extract_rtf_text(file_path: str) -> str:
    from striprtf.striprtf import rtf_to_text

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    return rtf_to_text(raw)


# ==========================================
# GENERATE KUIS TRAINING (dari isi dokumen)
# ==========================================
def generate_quiz_questions(content: str, num_questions: int = 5) -> list:
    """
    Minta Gemini bikin soal pilihan ganda dari isi dokumen SOP/training.
    Return list of dict: [{question, options:[4 opsi], correct_index}, ...]
    Format dipaksa JSON murni lewat instruksi eksplisit di prompt.
    """
    prompt = f"""Kamu adalah pembuat soal ujian untuk training karyawan baru.
Berdasarkan dokumen di bawah, buatkan {num_questions} soal pilihan ganda (4 opsi jawaban)
yang menguji pemahaman terhadap SOP/prosedur di dokumen ini.

ATURAN KETAT:
- Jawab HANYA dengan JSON murni, tanpa teks penjelasan apa pun di luar JSON.
- Jangan pakai markdown code block (tanpa ```json).
- Format persis seperti ini:
[
  {{"question": "...", "options": ["A", "B", "C", "D"], "correct_index": 0}},
  ...
]
- correct_index adalah index (0-3) dari jawaban yang benar di array "options".
- Soal harus berdasarkan isi dokumen, bukan pengetahuan umum di luar dokumen.

=== ISI DOKUMEN ===
{content[:6000]}
=== AKHIR DOKUMEN ===
"""
    response = _generate_with_fallback(prompt)
    raw = (response.text or "").strip()

    # Bersihkan kalau model tetap membungkus dengan markdown code block
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()

    import json

    questions = json.loads(raw)

    # Validasi minimal supaya tidak menyimpan data rusak ke database
    valid = []
    for q in questions:
        if (
            isinstance(q, dict)
            and "question" in q
            and "options" in q
            and len(q["options"]) == 4
            and "correct_index" in q
            and 0 <= q["correct_index"] <= 3
        ):
            valid.append(q)
    return valid


# ==========================================
# GENERATE DOKUMEN AI (draf, BUKAN dokumen resmi)
# ==========================================
DOC_TYPE_INSTRUCTIONS = {
    "SOP": "Format sebagai SOP: tujuan, ruang lingkup, lalu langkah-langkah bernomor yang jelas dan actionable.",
    "Form/Checklist": "Format sebagai daftar item/kolom isian singkat, cocok dijadikan tabel formulir.",
    "Surat": "Format sebagai surat resmi: kop (placeholder), tanggal, salam pembuka, isi, salam penutup, tempat tanda tangan.",
    "Kebijakan": "Format sebagai pernyataan kebijakan: latar belakang, ketentuan, sanksi/konsekuensi kalau relevan.",
    "Lainnya": "Format bebas namun tetap terstruktur dengan heading yang jelas.",
}


def generate_draft_document(
    topic: str, doc_type: str = "Lainnya", company_context: str = ""
) -> str:
    """
    Buat draf dokumen berdasarkan pengetahuan umum AI, disesuaikan jenis dokumen
    dan konteks perusahaan kalau ada. SELALU dilabeli sebagai draf.
    """
    instruction = DOC_TYPE_INSTRUCTIONS.get(
        doc_type, DOC_TYPE_INSTRUCTIONS["Lainnya"])
    prompt = f"""Kamu diminta membuat DRAF dokumen kerja untuk sebuah perusahaan.

Jenis dokumen: {doc_type}
Topik: {topic}
Konteks perusahaan (kalau ada): {company_context or "Tidak ada konteks tambahan."}

Instruksi format: {instruction}

ATURAN WAJIB:
- Di baris PALING ATAS, tulis persis: "[DRAF AI -- PERLU DIREVIEW SEBELUM DIGUNAKAN RESMI]"
- JANGAN mengarang detail teknis spesifik (angka, merek alat, dsb) yang tidak bisa
  dipastikan kebenarannya -- gunakan placeholder seperti "[isi sesuai SOP internal]".
"""
    response = _generate_with_fallback(prompt)
    return response.text or ""


def create_docx_bytes(title: str, content: str, logo_bytes: bytes = None) -> bytes:
    """Ubah teks jadi file .docx dari nol (tanpa template), opsional logo di kop."""
    from docx import Document
    from docx.shared import Inches
    import io

    doc = Document()
    if logo_bytes:
        try:
            doc.add_picture(io.BytesIO(logo_bytes), width=Inches(1.2))
        except Exception:
            pass  # logo korup/format tidak didukung -> lanjut tanpa logo, jangan gagalkan seluruh generate
    doc.add_heading(title, level=1)
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("[DRAF AI"):
            p = doc.add_paragraph()
            run = p.add_run(line)
            run.bold = True
        elif line.startswith("#"):
            doc.add_heading(line.lstrip("#").strip(), level=2)
        else:
            doc.add_paragraph(line)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def create_docx_from_template(template_bytes: bytes, title: str, content: str) -> bytes:
    """
    Sisipkan konten AI ke DALAM template .docx yang sudah diupload admin (berisi
    kop/logo/footer perusahaan) -- bukan bikin dokumen baru dari nol. Konten
    ditambahkan di akhir isi template, otomatis ikut font/style template.
    """
    from docx import Document
    import io

    doc = Document(io.BytesIO(template_bytes))
    doc.add_heading(title, level=1)
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("[DRAF AI"):
            p = doc.add_paragraph()
            run = p.add_run(line)
            run.bold = True
        elif line.startswith("#"):
            doc.add_heading(line.lstrip("#").strip(), level=2)
        else:
            doc.add_paragraph(line)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def create_pdf_bytes(title: str, content: str, logo_bytes: bytes = None) -> bytes:
    """Ubah teks jadi file .pdf sungguhan, opsional logo di kop."""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    import io

    styles = getSampleStyleSheet()
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
    )
    story = []
    if logo_bytes:
        try:
            story.append(Image(io.BytesIO(logo_bytes),
                         width=3 * cm, height=3 * cm))
            story.append(Spacer(1, 12))
        except Exception:
            pass
    story.append(Paragraph(title, styles["Title"]))
    story.append(Spacer(1, 16))
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 8))
            continue
        style = (
            styles["Heading2"]
            if line.startswith("[DRAF AI") or line.startswith("#")
            else styles["Normal"]
        )
        story.append(Paragraph(line.lstrip("#").strip(), style))
    doc.build(story)
    return buffer.getvalue()


def create_xlsx_bytes(title: str, rows: list) -> bytes:
    """
    Bikin file .xlsx dari list of dict (baris data). Dipakai untuk Form/Checklist
    kosong ATAU hasil kompilasi data nyata (daftar pelamar, rekap keuangan, dsb).
    """
    import openpyxl
    import io

    wb = openpyxl.Workbook()
    ws = wb.active
    # Excel batasi nama sheet 31 karakter
    ws.title = title[:31] if title else "Data"

    if rows:
        headers = list(rows[0].keys())
        ws.append(headers)
        for row in rows:
            ws.append([row.get(h, "") for h in headers])

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


# ==========================================
# GENERATE DARI DATA KOS (kompilasi data NYATA, bukan karangan AI)
# ==========================================
def determine_extraction_columns(user_request: str) -> list:
    """
    Tanya AI: berdasarkan permintaan user, kolom apa saja yang seharusnya ada
    di tabel hasil? Contoh: "daftar pelamar posisi marketing" -> ["Nama", "Posisi
    Dilamar", "Pendidikan", "Pengalaman (tahun)"].
    """
    prompt = f"""Permintaan pengguna: "{user_request}"

Tentukan kolom-kolom apa saja yang seharusnya ada di tabel hasil kompilasi data.
Jawab HANYA dengan JSON array of string, tanpa markdown code block, contoh:
["Nama", "Posisi Dilamar", "Pengalaman (tahun)"]
"""
    response = _generate_with_fallback(prompt)
    raw = (response.text or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    import json

    return json.loads(raw)


def extract_fields_from_document(
    content: str, columns: list, source_title: str
) -> dict:
    """
    Ekstrak nilai kolom yang diminta DARI ISI DOKUMEN NYATA (bukan mengarang).
    Kalau info tidak ada di dokumen, isi "-" -- JANGAN ditebak.
    """
    prompt = f"""Dokumen sumber ("{source_title}"):
{content[:4000]}

Ekstrak nilai untuk kolom berikut PERSIS dari isi dokumen di atas: {columns}
Kalau suatu info tidak disebutkan di dokumen, isi nilainya dengan "-" -- JANGAN menebak/mengarang.

Jawab HANYA dengan JSON object, tanpa markdown code block, contoh:
{{"Nama": "...", "Posisi Dilamar": "..."}}
"""
    response = _generate_with_fallback(prompt)
    raw = (response.text or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()
    import json

    return json.loads(raw)


# ==========================================
# ANALISIS DATA (AI cuma ekstrak kriteria, TIDAK menulis/menjalankan kode)
# ==========================================
def extract_analysis_criteria(question: str, columns: list) -> dict:
    """
    AI menerjemahkan pertanyaan bebas jadi kriteria terstruktur (JSON).
    Eksekusi filter dilakukan oleh kode Python kita sendiri (bukan AI) --
    supaya tidak ada risiko AI menjalankan kode sembarangan.

    Return dict:
    {
      "missing_info": null atau string pertanyaan klarifikasi,
      "filters": [{"column": "...", "operator": ">=|<=|==|contains", "value": "..."}],
      "sort_by": "nama_kolom atau null",
      "sort_desc": true/false
    }
    """
    prompt = f"""Kolom yang tersedia di data: {columns}

Pertanyaan pengguna: "{question}"

Terjemahkan pertanyaan ini jadi kriteria filter terstruktur. Jawab HANYA dengan JSON murni,
tanpa markdown code block, format persis:
{{
  "missing_info": null,
  "filters": [{{"column": "nama_kolom", "operator": ">=", "value": "1"}}],
  "sort_by": null,
  "sort_desc": false
}}

Operator yang valid HANYA: ">=", "<=", "==", "contains".
Kalau pertanyaan terlalu ambigu / kriteria penting belum disebutkan (misal user cuma
bilang "kehadiran bagus" tanpa angka), isi "missing_info" dengan pertanyaan klarifikasi
singkat dalam Bahasa Indonesia, dan biarkan "filters" jadi array kosong.
"""
    response = _generate_with_fallback(prompt)
    raw = (response.text or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
        raw = raw.strip()

    import json

    return json.loads(raw)
