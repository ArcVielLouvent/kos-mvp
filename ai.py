import os
import time
from google import genai
from google.genai import types
import streamlit as st


def get_client() -> genai.Client:
    """Inisialisasi Client menggunakan SDK google-genai yang mutakhir"""
    return genai.Client(api_key=st.secrets["GEMINI_API_KEY"])


def _call_with_retry(func, *args, max_retries: int = 4, base_delay: int = 3, **kwargs):
    """
    Coba ulang otomatis kalau kena error transient dari server Google
    (503 UNAVAILABLE / 429 rate limit) -- bukan bug di kode, murni server
    Google sedang sibuk. Backoff: 3s, 6s, 12s, 24s.
    """
    last_error = None
    for attempt in range(max_retries):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            last_error = e
            error_str = str(e)
            is_transient = any(
                code in error_str for code in ["503", "UNAVAILABLE", "429", "RESOURCE_EXHAUSTED"]
            )
            if is_transient and attempt < max_retries - 1:
                wait = base_delay * (2 ** attempt)
                time.sleep(wait)
                continue
            raise
    raise last_error


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list:
    """Sistem pemotong teks (Chunking) untuk mencegah Limit Token"""
    if not text:
        return []
    chunks = []
    i = 0
    while i < len(text):
        end = i + chunk_size
        if end < len(text):
            space_index = text.rfind(" ", i, end)
            if space_index > i:
                end = space_index
        chunks.append(text[i:end].strip())
        i = end - overlap
    return [c for c in chunks if len(c) > 10]


@st.cache_resource
def get_embedding_model() -> str:
    """Otomatis beralih ke lini Gemini Embedding terbaru"""
    return "gemini-embedding-2"


@st.cache_resource
def get_generation_model() -> str:
    """Menggunakan lini Gemini Flash paling mutakhir dan stabil"""
    return "gemini-3.5-flash"


def embed_text(text: str) -> list:
    """Sintaks ekstraksi array embedding dengan pembatasan dimensi ke 768"""
    client = get_client()
    model_name = get_embedding_model()

    result = _call_with_retry(
        client.models.embed_content,
        model=model_name,
        contents=text,
        config=types.EmbedContentConfig(output_dimensionality=768),
    )

    [embedding_obj] = result.embeddings
    return embedding_obj.values


def generate_answer(question: str, context_documents: list) -> str:
    client = get_client()
    context_parts = []
    for doc in context_documents:
        sumber = doc.get("metadata", {}).get("tipe_file", "Dokumen KOS")
        folder = doc.get("folder_path", "/")
        context_parts.append(
            f"Judul: {doc['title']}\nLokasi: {folder}\nTipe: {sumber}\nIsi:\n{doc['content']}"
        )

    context = "\n\n====================\n\n".join(context_parts)

    prompt = f"""Kamu adalah Knowledge Operating System (KOS), asisten cerdas internal perusahaan.
Jawab pertanyaan pengguna HANYA berdasarkan Dokumen Referensi di bawah ini.
Jika jawaban tidak ada di dalam dokumen referensi, katakan jujur bahwa informasi tersebut belum tersedia di database. Jangan pernah mengarang.

=== DOKUMEN REFERENSI ===
{context}
=== AKHIR DOKUMEN ===

Pertanyaan: {question}
"""
    model_name = get_generation_model()
    response = _call_with_retry(
        client.models.generate_content,
        model=model_name,
        contents=prompt,
    )
    return response.text


def extract_multimodal(file_path: str, mime_type: str, display_name: str) -> str:
    """
    Ekstraksi PDF, audio, dan video via Gemini File API (SDK google-genai).
    File diupload ke Google Cloud, dipoll sampai status ACTIVE (anti-halusinasi),
    lalu dihapus dari server setelah selesai.
    """
    client = get_client()

    try:
        uploaded_file = _call_with_retry(client.files.upload, file=file_path)

        # Polling anti-halusinasi: tunggu sampai file benar-benar siap diproses
        while uploaded_file.state.name == "PROCESSING":
            time.sleep(3)
            uploaded_file = client.files.get(name=uploaded_file.name)

        if uploaded_file.state.name == "FAILED":
            try:
                client.files.delete(name=uploaded_file.name)
            except Exception:
                pass
            raise ValueError(
                f"Google Gemini gagal memproses file '{display_name}'.")

        if "pdf" in mime_type:
            prompt = (
                "Baca seluruh dokumen PDF ini dengan saksama dan ekstrak seluruh "
                "teks serta tabel menjadi teks terstruktur murni."
            )
        elif "video" in mime_type:
            prompt = (
                "Tonton video ini dari awal hingga akhir. Buatkan transkrip sangat detail, "
                "termasuk kejadian visual, langkah-langkah teknis yang ditunjukkan, "
                "dan teks apa pun yang muncul di layar."
            )
        elif "audio" in mime_type:
            prompt = "Dengarkan audio ini dengan saksama dan buatkan transkrip teks yang lengkap dan akurat."
        elif "image" in mime_type:
            prompt = (
                "Amati gambar ini dengan saksama. Jika ada teks di dalamnya (dokumen yang "
                "difoto, poster, papan tulis, tangkapan layar), transkrip teksnya secara "
                "lengkap dan akurat. Jika ini foto biasa tanpa teks, deskripsikan isinya "
                "secara detail dan faktual."
            )
        else:
            prompt = "Ekstrak seluruh informasi dari file ini menjadi teks terstruktur murni."

        model_name = get_generation_model()
        response = _call_with_retry(
            client.models.generate_content,
            model=model_name,
            contents=[uploaded_file, prompt],
        )

        try:
            client.files.delete(name=uploaded_file.name)
        except Exception:
            pass

        hasil_teks = response.text
        if not hasil_teks or not hasil_teks.strip():
            raise ValueError(
                f"Tidak ada teks yang bisa diekstrak dari '{display_name}'.")

        return hasil_teks

    except Exception as e:
        raise RuntimeError(
            f"Gagal mengekstrak '{display_name}' via Google File API: {str(e)}")


# ==========================================
# EKSTRAKSI LOKAL (DOCX / PPTX / XLSX)
# Tidak lewat Gemini -- dokumen office tidak dipahami baik oleh document vision
# ==========================================
def extract_docx_text(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    return "\n".join(parts)


def extract_pptx_text(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        texts.append(text)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))

    return "\n\n".join(slides)


def extract_xlsx_text(file_path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    sheets = []

    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join(
                    str(c) if c is not None else "" for c in row))
        if rows:
            sheets.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))

    return "\n\n".join(sheets)


def extract_rtf_text(file_path: str) -> str:
    from striprtf.striprtf import rtf_to_text

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    return rtf_to_text(raw)
